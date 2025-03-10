import streamlit as st
from pptx import Presentation
import openai

def generate_ppt(title, content):
    prs = Presentation()
    
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    
    # Content Slides
    for i, point in enumerate(content, start=1):
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = f"Slide {i}: {point['heading']}"
        slide.shapes.placeholders[1].text = point['text']
    
    ppt_filename = "generated_presentation.pptx"
    prs.save(ppt_filename)
    return ppt_filename

def get_ppt_content(subject):
    prompt = f"Generate a 10-slide PowerPoint content about '{subject}'. Provide a heading and detailed content for each slide."
    response = openai.ChatCompletion.create(
        model="gpt-4",
        messages=[{"role": "system", "content": "You are an expert in creating PowerPoint presentations."},
                  {"role": "user", "content": prompt}]
    )
    
    slides_content = response["choices"][0]["message"]["content"].strip()
    
    # Parsing AI response
    slide_details = []
    slides = slides_content.split('\n\n')  # Assuming slides are separated by double newlines
    for slide in slides:
        if ':' in slide:
            heading, text = slide.split(':', 1)
            slide_details.append({"heading": heading.strip(), "text": text.strip()})
    
    return slide_details[:10]  # Ensure only 10 slides are used

# Streamlit UI
st.title("AI-Powered PPT Generator")
subject = st.text_input("Enter Subject Title")

if st.button("Generate PPT"):
    if subject:
        content = get_ppt_content(subject)
        ppt_file = generate_ppt(subject, content)
        with open(ppt_file, "rb") as f:
            st.download_button("Download PPT", f, file_name=ppt_file, mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
    else:
        st.error("Please enter a subject title.")
